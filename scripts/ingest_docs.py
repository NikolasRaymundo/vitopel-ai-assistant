# ingest_docs.py

# ——— Imports ———
import os
import fitz # PyMuPDF – PDF parsing (native & OCR fallback)
import pytesseract # OCR engine (Tesseract)
import pillow_heif # Adds HEIC support to Pillow
pillow_heif.register_heif_opener() # Enables PIL.Image.open() to read HEIC
from PIL import Image # Image handling for OCR
from docx import Document # DOCX parsing
import pandas as pd # Excel/CSV parsing
import json # Save processed outputs as JSON
import hashlib # Generate unique file hash
from pptx import Presentation # PPTX parsing
import zipfile # ZIP file extraction
import subprocess # Run LibreOffice conversions
import unicodedata # Normalize filenames (accents, symbols)
import shutil # Still potentially useful for file operations like os.remove
from datetime import datetime

# ——— Path Configuration ———
# Script is in 'scripts/', so project_root is one level up ('..').
project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))

RAW_FOLDER = os.path.join(project_root, "data", "raw_docs")
OUT_FOLDER = os.path.join(project_root, "data", "processed_texts")
MANIFEST_FILE = os.path.join(project_root, "data", "ingestion_manifest.json")

# ——— Global Constants ———
LEGACY_CONVERSION_MAP = {
    "ppt": "pptx",
    "doc": "docx",
    "xls": "xlsx"
}

# NOTE: Your individual parser functions (parse_pdf, parse_docx, etc.)
# should be defined AFTER this block, but BEFORE EXTENSION_PARSERS.
# We'll assume they are there.
# For EXTENSION_PARSERS to be defined here, the parser functions must be defined above it.
# So, the structure will be:
# ... (this block) ...
# ... (your parser functions: parse_pdf, parse_docx, etc.) ...
# ... EXTENSION_PARSERS = { ... } ...
# ... (other helper functions: converters, hash_file, normalize_filename, save_json) ...

# If your parser functions are already defined above this point in your original script,
# you can keep EXTENSION_PARSERS here. Otherwise, we'll define it later.
# For now, let's assume your parser functions are defined after this block.

# ——— Conversion helpers for legacy Office formats ———

def convert_ppt_to_pptx(path):
    try:
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pptx", path, "--outdir", os.path.dirname(path)],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL
        )
        converted_filename = os.path.splitext(os.path.basename(path))[0] + ".pptx"
        converted_path = os.path.join(os.path.dirname(path), converted_filename)
        return converted_path if os.path.exists(converted_path) else None
    except Exception as e:
        print(f"[fail] PPT conversion error: {path}: {e}")
        return None

def convert_doc_to_docx(path):
    try:
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "docx", path, "--outdir", os.path.dirname(path)],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL
        )
        converted_filename = os.path.splitext(os.path.basename(path))[0] + ".docx"
        converted_path = os.path.join(os.path.dirname(path), converted_filename)
        return converted_path if os.path.exists(converted_path) else None
    except Exception as e:
        print(f"[fail] DOC conversion error: {path}: {e}")
        return None

def convert_xls_to_xlsx(path):
    try:
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "xlsx", path, "--outdir", os.path.dirname(path)],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL
        )
        converted_filename = os.path.splitext(os.path.basename(path))[0] + ".xlsx"
        converted_path = os.path.join(os.path.dirname(path), converted_filename)
        return converted_path if os.path.exists(converted_path) else None
    except Exception as e:
        print(f"[fail] XLS conversion error: {path}: {e}")
        return None

# ——— Utility: File hasher ———
def hash_file(filepath):
    with open(filepath, "rb") as f:
        return hashlib.sha256(f.read()).hexdigest()

# ——— Parsers for different file types ———

def parse_pdf(path):
    text = ""
    try:
        doc = fitz.open(path)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page_text = page.get_text("text")
            if page_text and page_text.strip():
                text += page_text + "\n"
        if text.strip():
            return text.strip()
        else:
            print(f"[info] No text layer in PDF {os.path.basename(path)}, attempting OCR...")
            text_ocr = ""
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(dpi=300)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                ocr_text_page = pytesseract.image_to_string(img, lang="por+eng")
                if ocr_text_page and ocr_text_page.strip():
                    text_ocr += ocr_text_page + "\n"
            return text_ocr.strip() if text_ocr.strip() else None
    except Exception as e:
        print(f"[PDF/OCR FAIL] {os.path.basename(path)}: {e}")
        return None

def parse_docx(path):
    try:
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text and p.text.strip())
    except Exception as e:
        print(f"[DOCX FAIL] {os.path.basename(path)}: {e}")
        return None

def parse_excel(path):
    try:
        df = pd.read_excel(path, engine='openpyxl')
        return df.to_string(index=False)
    except Exception as e:
        print(f"[XLSX FAIL] {os.path.basename(path)}: {e}")
        return None

def parse_csv(path):
    try:
        df = pd.read_csv(path)
        return df.to_string(index=False)
    except Exception as e:
        print(f"[CSV FAIL] {os.path.basename(path)}: {e}")
        return None

def parse_image(path):
    try:
        img = Image.open(path)
        print(f"[debug] Opened image: {path}, format: {img.format}, mode: {img.mode}")
        if img.format and img.format.lower() in ["heif", "heic", "webp"]:
            img = img.convert("RGB")
        text = pytesseract.image_to_string(img, lang="por+eng")
        return text.strip() if text.strip() else None
    except Exception as e:
        print(f"[IMG FAIL] {os.path.basename(path)}: {type(e).__name__}: {e}")
        return None

def parse_txt(path):
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read().strip()
    except Exception as e:
        print(f"[TXT FAIL] {os.path.basename(path)}: {e}")
        return None

def parse_pptx(path):
    try:
        prs = Presentation(path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    text_parts.append(shape.text.strip())
        return "\n".join(text_parts) if text_parts else None
    except Exception as e:
        print(f"[PPTX FAIL] {os.path.basename(path)}: {e}")
        return None

# --- Registry that maps extensions to handlers ---
# Ensure all functions listed here (parse_pdf, parse_docx, etc.)
# are defined above this point in your script.
EXTENSION_PARSERS = {
    "pdf": parse_pdf,
    "docx": parse_docx,
    "xlsx": parse_excel,
    "csv": parse_csv,
    "png": parse_image,
    "jpg": parse_image,
    "jpeg": parse_image,
    "heic": parse_image,
    "tif": parse_image,
    "tiff": parse_image,
    "txt": parse_txt,
    "pptx": parse_pptx,
}

# ——— Save parsed output as .json ———
def save_json(filename_without_json_ext, data):
    out_path = os.path.join(OUT_FOLDER, filename_without_json_ext + ".json")
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)

# ——— Normalize weird filenames from ZIP extractions ———
def normalize_filename(name):
    return unicodedata.normalize("NFKD", name).encode("ascii", "ignore").decode("ascii")

# --- Manifest Handling Functions ---
def load_manifest():
    """Loads the ingestion manifest if it exists and is valid JSON."""
    if os.path.exists(MANIFEST_FILE):
        try:
            with open(MANIFEST_FILE, 'r', encoding='utf-8') as f:
                manifest_data = json.load(f)
                print(f"[info] Loaded existing manifest from {MANIFEST_FILE}")
                return manifest_data
        except json.JSONDecodeError:
            print(f"[warn] Manifest file {MANIFEST_FILE} is corrupted or not valid JSON. Starting with a new manifest.")
            return {} # Return empty dict if corrupted
        except Exception as e:
            print(f"[warn] Could not load manifest file {MANIFEST_FILE} due to an error: {e}. Starting with a new manifest.")
            return {}
    else:
        print(f"[info] No manifest file found at {MANIFEST_FILE}. A new one will be created.")
        return {} # Return empty dict if no file

def save_manifest(manifest_data):
    """Saves the provided manifest data to the manifest file."""
    try:
        # Ensure the 'data' directory exists for the manifest file
        os.makedirs(os.path.dirname(MANIFEST_FILE), exist_ok=True) 
        with open(MANIFEST_FILE, 'w', encoding='utf-8') as f:
            json.dump(manifest_data, f, indent=4, ensure_ascii=False) # Use indent=4 for readability
        print(f"[info] Manifest saved to {MANIFEST_FILE}")
    except Exception as e:
        print(f"[error] Could not save manifest to {MANIFEST_FILE}: {e}")

# === PHASE-SPECIFIC FUNCTIONS ===

def phase_a_unzip_archives(base_raw_folder):
    """
    Scans for ZIP files within base_raw_folder (and its subdirectories),
    extracts their contents into a new subfolder named after the ZIP file
    (minus .zip extension), in the same location as the ZIP.
    Filenames from within the ZIP are normalized using the project's 
    normalize_filename function.
    It attempts to clean up __MACOSX folders and flatten single-folder ZIP structures.
    Returns a list of full paths to the original ZIP files that were processed.
    """
    print("--- Starting Phase A: Unzipping Archives ---")
    processed_zip_fpaths = [] 

    for root, dirs, files in os.walk(base_raw_folder):
        for filename in files:
            if filename.lower().endswith(".zip"):
                zip_fpath = os.path.join(root, filename)
                print(f"[zip] Found archive: {zip_fpath}")
                
                processed_zip_fpaths.append(zip_fpath) 
                
                zip_name_without_ext = os.path.splitext(filename)[0]
                # Normalize the folder name created from the ZIP filename itself
                normalized_zip_folder_name = normalize_filename(zip_name_without_ext)
                if not normalized_zip_folder_name.strip(): # Handle if zip name becomes empty after normalization
                    normalized_zip_folder_name = "untitled_zip_extraction"
                    print(f"    [warn] ZIP filename '{zip_name_without_ext}' normalized to empty or whitespace, using '{normalized_zip_folder_name}' for extraction folder.")

                extraction_target_folder = os.path.join(root, normalized_zip_folder_name)
                
                extraction_succeeded_at_all = False
                try:
                    if not os.path.exists(extraction_target_folder):
                        os.makedirs(extraction_target_folder)
                        print(f"    Created extraction folder: {extraction_target_folder}")
                    else:
                        print(f"    [warn] Extraction folder '{extraction_target_folder}' already exists. Files may be overwritten if names collide after normalization.")

                    with zipfile.ZipFile(zip_fpath, 'r') as zip_ref:
                        for member in zip_ref.infolist():
                            if member.is_dir():
                                continue # Skip directory entries themselves
                            
                            # Skip common Mac resource files and other unwanted files
                            if member.filename.startswith('__MACOSX/') or \
                               os.path.basename(member.filename) == '.DS_Store' or \
                               not member.filename.strip(): # Skip empty filenames
                                print(f"    [skip] Ignored metadata or empty entry: {member.filename}")
                                continue

                            # Construct the target path for the member, normalizing each path component
                            path_parts = member.filename.split('/')
                            normalized_path_parts = []
                            for part in path_parts:
                                if part: # Ensure part is not empty
                                    normalized_part = normalize_filename(part)
                                    if normalized_part.strip(): # Ensure normalized part is not empty
                                        normalized_path_parts.append(normalized_part)
                                    else:
                                        # If a path part becomes empty after normalization, this could be problematic
                                        # For now, we'll just note it and it might get skipped if normalized_path_parts ends up empty
                                        print(f"        [warn] Path component '{part}' from '{member.filename}' normalized to empty string.")
                            
                            if not normalized_path_parts:
                                print(f"    [skip] Member '{member.filename}' resulted in an empty path after normalization. Skipping.")
                                continue
                            
                            # Create the on-disk path using normalized components
                            # Example: if member.filename is "SubArchive/4 Módulo Utilidades.pptx"
                            # normalized_path_parts could be ["SubArchive", "4 Modulo Utilidades.pptx"]
                            # target_fpath_on_disk will be extraction_target_folder/SubArchive/4 Modulo Utilidades.pptx
                            target_fpath_on_disk = os.path.join(extraction_target_folder, *normalized_path_parts)
                            
                            # Ensure parent directory for the file exists
                            target_member_parent_dir = os.path.dirname(target_fpath_on_disk)
                            if not os.path.exists(target_member_parent_dir):
                                os.makedirs(target_member_parent_dir)

                            # Extract file data to the normalized path
                            try:
                                with zip_ref.open(member) as source_file, open(target_fpath_on_disk, "wb") as target_file:
                                    shutil.copyfileobj(source_file, target_file)
                                # print(f"    Extracted '{member.filename}' to '{target_fpath_on_disk}'") # Can be verbose
                                extraction_succeeded_at_all = True 
                            except Exception as e_extract_member:
                                print(f"    [fail] Could not extract member '{member.filename}' to '{target_fpath_on_disk}': {e_extract_member}")
                    
                    if extraction_succeeded_at_all:
                        print(f"    Successfully extracted and normalized contents of '{filename}' to '{extraction_target_folder}'")
                    else:
                        print(f"    [warn] No files were successfully extracted from '{filename}'.")


                except zipfile.BadZipFile:
                    print(f"    [fail] Bad ZIP file: '{filename}'. Contents not extracted.")
                except Exception as e: # Other errors like permission issues creating folders
                    print(f"    [fail] Could not process ZIP '{filename}': {e}")

                # --- Post-extraction cleanup for single root folder (still useful) ---
                if extraction_succeeded_at_all and os.path.exists(extraction_target_folder): # Only if extraction folder was created
                    try:
                        current_contents = os.listdir(extraction_target_folder)
                        # Filter out macOS specific hidden folders if any slipped through normalization or were top-level in zip
                        significant_contents = [item for item in current_contents if not item.startswith('.') and item.lower() != '__macosx']

                        if len(significant_contents) == 1:
                            single_item_name = significant_contents[0]
                            path_to_single_item = os.path.join(extraction_target_folder, single_item_name)
                            
                            if os.path.isdir(path_to_single_item):
                                print(f"    [info] Extraction appears to have a single root folder: '{single_item_name}'. Promoting its contents.")
                                # Define a temporary directory name for promotion to avoid conflicts
                                temp_promote_dir = os.path.join(extraction_target_folder, single_item_name + "_promotetemp")
                                
                                try:
                                    os.rename(path_to_single_item, temp_promote_dir) # Rename single item dir to temp name
                                    
                                    items_moved_successfully = True
                                    for item_to_move in os.listdir(temp_promote_dir):
                                        source_item_path = os.path.join(temp_promote_dir, item_to_move)
                                        dest_item_path = os.path.join(extraction_target_folder, item_to_move)
                                        
                                        if os.path.exists(dest_item_path): # Collision check
                                            print(f"        [warn] Item '{item_to_move}' already exists in target '{extraction_target_folder}'. Attempting to remove existing and replace.")
                                            try:
                                                if os.path.isfile(dest_item_path): os.remove(dest_item_path)
                                                elif os.path.isdir(dest_item_path): shutil.rmtree(dest_item_path)
                                                shutil.move(source_item_path, dest_item_path)
                                            except Exception as e_replace:
                                                print(f"        [fail] Could not replace '{dest_item_path}'. Item '{item_to_move}' not promoted. {e_replace}")
                                                items_moved_successfully = False
                                        else:
                                            shutil.move(source_item_path, dest_item_path)
                                    
                                    if items_moved_successfully:
                                        shutil.rmtree(temp_promote_dir) # Remove the (now empty) temporary promotion directory
                                        print(f"    [cleanup] Successfully promoted contents from '{single_item_name}'.")
                                    else: # If not all items moved, move the temp dir back
                                        os.rename(temp_promote_dir, path_to_single_item)
                                        print(f"    [warn] Not all items from '{single_item_name}' were promoted due to conflicts. Original structure within '{single_item_name}' restored from temp.")

                                except Exception as e_promote:
                                    print(f"    [warn] Error during content promotion for '{single_item_name}': {e_promote}. Structure may be as originally extracted.")
                                    # If temp_promote_dir still exists and original path_to_single_item does not, rename back
                                    if os.path.exists(temp_promote_dir) and not os.path.exists(path_to_single_item):
                                        try:
                                            os.rename(temp_promote_dir, path_to_single_item)
                                        except Exception as e_rename_back:
                                            print(f"        [warn] Could not rename temp promotion directory back: {e_rename_back}")


                    except Exception as e_listdir:
                        print(f"    [warn] Could not inspect contents of '{extraction_target_folder}' for single root folder heuristic: {e_listdir}")
    
    print(f"--- Finished Phase A: Attempted to process {len(processed_zip_fpaths)} ZIP archives ---")
    return processed_zip_fpaths

def phase_b_convert_legacy_files(base_raw_folder, legacy_map):
    """
    Scans for legacy Office files within base_raw_folder (and its subdirectories).
    If a modern equivalent already exists, the legacy file is marked for deletion.
    Otherwise, it attempts to convert legacy files to modern formats in the same location.
    If conversion is successful, the original legacy file is marked for deletion.
    Returns a list of full paths to the original legacy files that should be deleted.
    """
    print("--- Starting Phase B: Converting/Checking Legacy Files ---")
    legacy_fpaths_to_delete = [] 

    current_files_in_raw = []
    for root, dirs, files in os.walk(base_raw_folder):
        for filename in files:
            current_files_in_raw.append(os.path.join(root, filename))

    for original_fpath in current_files_in_raw:
        # Check if file still exists, as it might have been a legacy file already converted
        # from a previous entry in current_files_in_raw (e.g. if scan order was not ideal
        # or if a .doc and .docx existed and .doc was processed creating another .docx).
        # This is less likely now with the more direct logic but a safe check.
        if not os.path.isfile(original_fpath):
            continue

        filename = os.path.basename(original_fpath)
        root = os.path.dirname(original_fpath) # Get the correct root for this file
        current_ext = os.path.splitext(filename)[1].lower().lstrip(".")
        
        if current_ext in legacy_map:
            base_name_no_ext = os.path.splitext(filename)[0]
            modern_ext = legacy_map[current_ext]
            expected_modern_filename = base_name_no_ext + "." + modern_ext
            expected_modern_fpath = os.path.join(root, expected_modern_filename)

            if os.path.exists(expected_modern_fpath):
                # Modern version already exists. Mark legacy file for deletion.
                # We don't need to compare timestamps for the deletion decision.
                # If the user wants the freshest conversion, they should ensure no modern version exists first,
                # or we could add logic to overwrite if legacy is newer (but that makes it more complex again).
                # For "cheeky deletion", if modern exists, legacy is redundant.
                print(f"    [info] Modern version '{expected_modern_filename}' already exists for legacy file '{filename}'. Marking legacy for deletion.")
                if original_fpath not in legacy_fpaths_to_delete: # Avoid duplicates if somehow scanned twice
                     legacy_fpaths_to_delete.append(original_fpath)
            else:
                # Modern version does not exist, so attempt conversion.
                print(f"[convert] Attempting to convert: {original_fpath}")
                converted_modern_fpath_output = None # Path returned by converter
                if current_ext == "doc":
                    converted_modern_fpath_output = convert_doc_to_docx(original_fpath)
                elif current_ext == "ppt":
                    converted_modern_fpath_output = convert_ppt_to_pptx(original_fpath)
                elif current_ext == "xls":
                    converted_modern_fpath_output = convert_xls_to_xlsx(original_fpath)
                
                # Check if the expected modern file was indeed created by the conversion
                if os.path.exists(expected_modern_fpath):
                    print(f"    -> Successfully converted '{filename}' to '{expected_modern_filename}'")
                    if original_fpath not in legacy_fpaths_to_delete:
                        legacy_fpaths_to_delete.append(original_fpath)
                else:
                    print(f"    [fail] Conversion failed or expected modern file '{expected_modern_filename}' not found for '{filename}'. Original file will be kept.")
        
    print(f"--- Finished Phase B: Marked {len(legacy_fpaths_to_delete)} original legacy files for deletion ---")
    return legacy_fpaths_to_delete

def phase_c_normalize_all_item_names_in_folder(folder_path):
    """
    Recursively scans all files and directories within folder_path and renames
    them to their normalized (ASCII) versions if they differ.
    Handles potential name collisions during renaming by appending a counter.
    Operates from deepest items upwards (topdown=False).
    """
    print("--- Starting Sub-Phase of C: Normalizing all item names in RAW_FOLDER to ASCII ---")
    renamed_count = 0
    # Walk from deepest to shallowest (topdown=False) to handle renaming of directory contents before the directory itself
    for root, dirs, files in os.walk(folder_path, topdown=False):
        # Normalize filenames first
        for name in files:
            original_fpath = os.path.join(root, name)
            
            # Ensure name is a string (os.walk can sometimes yield bytes on some systems/Python versions)
            current_name_str = name
            if isinstance(name, bytes):
                try:
                    current_name_str = name.decode(errors='surrogateescape') # Try to decode gracefully
                except UnicodeDecodeError:
                    print(f"    [warn] Could not decode filename bytes: {name}. Skipping normalization for this item.")
                    continue
            
            normalized_name = normalize_filename(current_name_str) # Your existing function

            if current_name_str != normalized_name:
                normalized_fpath_base = os.path.join(root, normalized_name)
                final_normalized_fpath = normalized_fpath_base
                
                counter = 1
                # Check for collision with existing files/dirs for the new normalized name
                while os.path.exists(final_normalized_fpath):
                    if os.path.normcase(original_fpath) == os.path.normcase(final_normalized_fpath):
                         print(f"    [info] File '{current_name_str}' normalized to '{normalized_name}' which is effectively the same path. Skipping explicit rename.")
                         current_name_str = normalized_name # Update current_name_str to avoid rename attempt
                         break 

                    base, ext = os.path.splitext(normalized_name)
                    final_normalized_fpath = os.path.join(root, f"{base}_{counter}{ext}")
                    counter += 1
                
                if current_name_str != normalized_name: 
                    try:
                        os.rename(original_fpath, final_normalized_fpath)
                        print(f"    [rename] Normalized file: '{original_fpath}' to '{final_normalized_fpath}'")
                        renamed_count +=1
                    except Exception as e:
                        print(f"    [rename_fail] Could not normalize file '{original_fpath}' to '{final_normalized_fpath}': {e}")
        
        # Normalize directory names second
        for name in dirs:
            original_dpath = os.path.join(root, name)
            current_name_str = name
            if isinstance(name, bytes):
                try:
                    current_name_str = name.decode(errors='surrogateescape')
                except UnicodeDecodeError:
                    print(f"    [warn] Could not decode dirname bytes: {name}. Skipping normalization for this item.")
                    continue

            normalized_name = normalize_filename(current_name_str)

            if current_name_str != normalized_name:
                normalized_dpath_base = os.path.join(root, normalized_name)
                final_normalized_dpath = normalized_dpath_base

                counter = 1
                while os.path.exists(final_normalized_dpath):
                    if os.path.normcase(original_dpath) == os.path.normcase(final_normalized_dpath):
                        print(f"    [info] Directory '{current_name_str}' normalized to '{normalized_name}' which is effectively the same path. Skipping explicit rename.")
                        current_name_str = normalized_name
                        break
                    
                    final_normalized_dpath = os.path.join(root, f"{normalized_name}_{counter}")
                    counter += 1
                
                if current_name_str != normalized_name:
                    try:
                        os.rename(original_dpath, final_normalized_dpath)
                        print(f"    [rename] Normalized directory: '{original_dpath}' to '{final_normalized_dpath}'")
                        renamed_count +=1
                    except Exception as e:
                        print(f"    [rename_fail] Could not normalize directory '{original_dpath}' to '{final_normalized_dpath}': {e}")
    
    if renamed_count > 0:
        print(f"--- Finished Sub-Phase of C: Normalized {renamed_count} item names in RAW_FOLDER ---")
    else:
        print("--- Sub-Phase of C: No item names in RAW_FOLDER required normalization ---")

def get_normalized_path_as_filename_base(full_fpath, base_folder_path):
    """
    Calculates a relative path and normalizes its components to create a single
    string suitable for a filename base.
    Example: 
    full_fpath = '/path/to/project/data/raw_docs/My Folder/My File (new).docx'
    base_folder_path = '/path/to/project/data/raw_docs'
    Returns: 'My Folder_My File (new).docx' (after normalization of components)
    """
    if not full_fpath.startswith(base_folder_path):
        # Should not happen if called correctly, fallback to just basename
        return normalize_filename(os.path.basename(full_fpath))

    relative_path = os.path.relpath(full_fpath, base_folder_path)
    
    parts = []
    head = relative_path
    while True:
        head, tail = os.path.split(head)
        if tail: # tail is a path component (file or dir name)
            parts.insert(0, normalize_filename(tail))
        else: # No tail means we've reached the top of the relative path
            if head: # If there's still a head (e.g. on Unix for /foo from /)
                parts.insert(0, normalize_filename(head))
            break
        if not head: # Stop if head is empty (processed all parts)
            break
            
    return "_".join(parts)

def process_modern_file_to_json(fpath_in_raw, manifest_key, current_file_hash, out_folder_path_flat): # Note: signature unchanged from last version
    """
    Processes a single modern file: parses its content and saves it to a JSON file
    in the flat out_folder_path_flat. Uses the provided manifest_key and current_file_hash.
    The JSON filename is derived from the manifest_key and hash.
    Returns the leaf filename of the created JSON file if successful, None otherwise.
    """
    try:
        # fpath_in_raw is already an ASCII path because of Phase C
        # manifest_key is already the normalized path identifier (e.g., Folder_File.ext)
        # current_file_hash is the hash of fpath_in_raw

        current_filename_on_disk = os.path.basename(fpath_in_raw) # This is already ASCII
        current_ext = os.path.splitext(current_filename_on_disk)[1].lower().lstrip(".")

        # 1. Get the parser for this modern file type
        parser = EXTENSION_PARSERS.get(current_ext)
        if not parser:
            print(f"    [skip] No parser available for modern file type: '{current_ext}' for file '{current_filename_on_disk}'")
            return None # Return None on failure

        # 2. Parse the file content
        parsed_text = parser(fpath_in_raw)
        if parsed_text is None: 
            print(f"    [fail] Parsing returned no content for: '{current_filename_on_disk}'")
            return None # Return None on failure
        
        # 3. Construct the output JSON filename (using provided manifest_key and current_file_hash)
        output_json_leafname_base = manifest_key + "." + current_file_hash[:8]
        actual_output_json_filename = output_json_leafname_base + ".json"
        
        # 4. Prepare data to save
        data_to_save = {
            "file_name": current_filename_on_disk,  # Basename from RAW_FOLDER (now ASCII)
            "file_type": current_ext,
            "text": parsed_text,
            "hash": current_file_hash, # Use the provided hash
            "source_path_in_raw": fpath_in_raw, # Full path in RAW_FOLDER (now ASCII)
            "normalized_identifier": manifest_key, # Use the provided manifest_key
            "ingestion_timestamp": datetime.now().isoformat()
        }
        
        # 5. Save the JSON file 
        # (Your existing save_json function prepends OUT_FOLDER and appends .json)
        save_json(output_json_leafname_base, data_to_save) 
        
        print(f"    [ok] Processed and saved to: {actual_output_json_filename}")
        return actual_output_json_filename # Return the leaf filename of the created JSON

    except Exception as e:
        print(f"    [error] Failed to process modern file '{os.path.basename(fpath_in_raw)}': {e}")
        return None # Return None on failure
    
def cleanup_orphaned_json_files(out_folder_path, active_manifest):
    """
    Deletes JSON files from out_folder_path that are not listed as 
    'output_json_filename' in the active_manifest.
    """
    print("--- Starting Cleanup of Orphaned JSON Files ---")
    if not active_manifest:
        print("    [info] Manifest is empty. No cleanup performed (or all files would be considered orphaned).")
        return

    active_json_filenames = set()
    for key, entry_data in active_manifest.items():
        if isinstance(entry_data, dict) and "output_json_filename" in entry_data:
            active_json_filenames.add(entry_data["output_json_filename"])
        else:
            print(f"    [warn] Manifest entry for key '{key}' is missing 'output_json_filename' or is not a dictionary. Value: {entry_data}")


    if not active_json_filenames:
        print("    [info] No active JSON filenames found in manifest. No cleanup performed.")
        return

    deleted_count = 0
    found_count = 0
    try:
        for filename in os.listdir(out_folder_path):
            if filename.endswith(".json"): # Only consider JSON files
                found_count += 1
                if filename not in active_json_filenames:
                    fpath_to_delete = os.path.join(out_folder_path, filename)
                    try:
                        os.remove(fpath_to_delete)
                        print(f"    [cleanup] Deleted orphaned JSON: {filename}")
                        deleted_count += 1
                    except Exception as e:
                        print(f"    [warn] Could not delete orphaned JSON '{filename}': {e}")
    except FileNotFoundError:
        print(f"    [warn] Output folder {out_folder_path} not found during cleanup. No cleanup performed.")
    except Exception as e:
        print(f"    [error] An error occurred listing files for cleanup in {out_folder_path}: {e}")


    print(f"--- Finished Cleanup: Found {found_count} JSON files, deleted {deleted_count} orphaned files. ---")

# === MAIN ORCHESTRATOR ===
def main():
    print("🚀 Starting Vitopel Document Ingestion Process...")

    # Ensure base folders exist
    os.makedirs(RAW_FOLDER, exist_ok=True)
    os.makedirs(OUT_FOLDER, exist_ok=True)

    # Load existing manifest at the start
    manifest = load_manifest()
    new_manifest_for_this_run = {} # To build the manifest for this current run

    # --- Phase A: Unzip Archives ---
    processed_zip_archive_fpaths = phase_a_unzip_archives(RAW_FOLDER)

    # --- Phase B: Convert Legacy Files ---
    legacy_fpaths_marked_for_deletion = phase_b_convert_legacy_files(RAW_FOLDER, LEGACY_CONVERSION_MAP)
    
    # --- Phase C: Delete Original Legacy Files, Processed ZIP Archives, and Normalize All Names ---
    print("--- Starting Phase C: Deleting Processed Original Files & Normalizing All Item Names ---")
    deleted_legacy_count = 0
    for fpath in legacy_fpaths_marked_for_deletion: 
        try:
            if os.path.exists(fpath): 
                os.remove(fpath)
                print(f"[cleanup] Deleted original legacy file: {fpath}")
                deleted_legacy_count += 1
        except Exception as e:
            print(f"[warn] Could not delete original legacy file {fpath}: {e}")
    if deleted_legacy_count > 0:
        print(f"Cleaned up {deleted_legacy_count} original legacy files.")
    
    deleted_zip_count = 0
    for fpath in processed_zip_archive_fpaths:
        zip_original_basename = os.path.basename(fpath)
        zip_name_without_ext = os.path.splitext(zip_original_basename)[0]
        normalized_zip_folder_name_for_check = normalize_filename(zip_name_without_ext)
        if not normalized_zip_folder_name_for_check.strip():
             normalized_zip_folder_name_for_check = "untitled_zip_extraction"

        extraction_folder_for_this_zip = os.path.join(os.path.dirname(fpath), normalized_zip_folder_name_for_check)
        
        if os.path.isdir(extraction_folder_for_this_zip): 
            try:
                if os.path.exists(fpath): 
                    os.remove(fpath)
                    print(f"[cleanup] Deleted original ZIP archive: {fpath}")
                    deleted_zip_count +=1
            except Exception as e:
                print(f"[warn] Could not delete original ZIP {fpath}: {e}")
        else:
            print(f"[info] Did not delete ZIP {fpath}. Its expected normalized extraction folder '{extraction_folder_for_this_zip}' was not found (extraction might have failed or was incomplete).")
    if deleted_zip_count > 0:
        print(f"Cleaned up {deleted_zip_count} original ZIP archives.")

    phase_c_normalize_all_item_names_in_folder(RAW_FOLDER)
    print("--- Finished Phase C ---")

    # --- Phase D: Process all (now modern and ASCII-named) files in RAW_FOLDER to JSONs ---
    print("--- Starting Phase D: Processing Modern Files to JSON (with Manifest Check) ---")
    
    files_to_consider_for_json = []
    for root, dirs, files in os.walk(RAW_FOLDER):
        for filename in files:
            if filename.lower().endswith(".zip"): 
                continue 
            
            ext = os.path.splitext(filename)[1].lower().lstrip(".")
            if ext in EXTENSION_PARSERS: 
                 file_path = os.path.join(root, filename)
                 files_to_consider_for_json.append(file_path)
            else:
                if ext not in LEGACY_CONVERSION_MAP: 
                    print(f"    [info] Skipping file with unknown/unhandled extension in Phase D: {filename} at {root}")

    if not files_to_consider_for_json:
        print("No files found in RAW_FOLDER for JSON processing after cleanup.")
    else:
        print(f"Found {len(files_to_consider_for_json)} files to potentially process into JSON.")

    processed_this_run_count = 0
    skipped_count = 0
    failed_count = 0
    
    for i, fpath_in_raw in enumerate(files_to_consider_for_json):
        print(f"\n  Considering file ({i+1}/{len(files_to_consider_for_json)}): {fpath_in_raw}...")
        
        try:
            manifest_key = get_normalized_path_as_filename_base(fpath_in_raw, RAW_FOLDER)
            current_file_hash = hash_file(fpath_in_raw)

            # Manifest Check
            if manifest_key in manifest and manifest[manifest_key].get("source_file_hash") == current_file_hash:
                print(f"    [skip] File '{os.path.basename(fpath_in_raw)}' is unchanged (manifest). Output: {manifest[manifest_key].get('output_json_filename')}")
                new_manifest_for_this_run[manifest_key] = manifest[manifest_key] # Carry over old entry
                skipped_count += 1
                continue # Skip to next file
            
            # If not skipped, then process
            print(f"    [process] New or changed file: {os.path.basename(fpath_in_raw)}")
            output_json_leafname = process_modern_file_to_json(fpath_in_raw, manifest_key, current_file_hash, OUT_FOLDER)
            
            if output_json_leafname:
                new_manifest_for_this_run[manifest_key] = {
                    "source_file_hash": current_file_hash,
                    "output_json_filename": output_json_leafname, # Leaf name returned by processor
                    "processed_at": datetime.now().isoformat() 
                    # Add status if needed: "status": "processed"
                }
                processed_this_run_count += 1
            else:
                # process_modern_file_to_json already prints error, just count failure
                failed_count += 1
                # Optionally, add an error entry to manifest, or exclude it
                # new_manifest_for_this_run[manifest_key] = {"source_file_hash": current_file_hash, "status": "processing_error", "processed_at": datetime.now().isoformat()}


        except FileNotFoundError:
            print(f"    [error] File not found during Phase D processing (may have been deleted unexpectedly): {fpath_in_raw}")
            failed_count += 1
        except Exception as e:
            print(f"    [error] Unexpected error during Phase D processing for {fpath_in_raw}: {e}")
            failed_count += 1
        
    # Save the new manifest for this run (overwrites old one)
    save_manifest(new_manifest_for_this_run)
    # Call cleanup immediately after saving the definitive manifest for this run
    cleanup_orphaned_json_files(OUT_FOLDER, new_manifest_for_this_run)

    print(f"\n--- Finished Phase D ---")
    print(f"  Files processed this run: {processed_this_run_count}")
    print(f"  Files skipped (unchanged): {skipped_count}")
    if failed_count > 0:
        print(f"  Files failed processing: {failed_count}")
        
    print(f"\n🎯 Ingestion Script Fully Executed.")

if __name__ == "__main__":
    main()